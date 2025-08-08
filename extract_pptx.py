#!/usr/bin/env python3
"""Extract text content from Terra35.pptx for analysis."""

from pptx import Presentation
import sys

def extract_pptx_content(filename):
    """Extract all text content from a PowerPoint presentation."""
    prs = Presentation(filename)
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*60}")
        print(f"SLIDE {i}")
        print('='*60)
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(text)
                    print('-'*40)
        
        # Extract notes if present
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print("\n[SPEAKER NOTES:]")
                print(notes)

if __name__ == "__main__":
    extract_pptx_content("Terra35.pptx")